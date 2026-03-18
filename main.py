from fastapi import FastAPI, UploadFile, File, HTTPException, Request, Body
from fastapi.responses import FileResponse, HTMLResponse, StreamingResponse
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from fastapi.templating import Jinja2Templates

import os
import uuid
import subprocess
import io
import base64
from PIL import Image

# Conversion libraries
from pdf2docx import Converter
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches

# ------------------ APP ------------------
app = FastAPI()

# ------------------ CORS ------------------
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

# ------------------ FOLDERS ------------------
UPLOAD_DIR = "uploads"
OUTPUT_DIR = "outputs"

os.makedirs(UPLOAD_DIR, exist_ok=True)
os.makedirs(OUTPUT_DIR, exist_ok=True)

# ------------------ STATIC & TEMPLATES ------------------
templates = Jinja2Templates(directory="templates")
app.mount("/static", StaticFiles(directory="static"), name="static")

# ------------------ PATHS ------------------
LIBREOFFICE_PATH = r"C:\Program Files\LibreOffice\program\soffice.exe"
POPPLER_PATH = r"C:\Users\veer\Downloads\Release-25.12.0-0\Library\bin"

# =========================================================
# UI ROUTES
# =========================================================

@app.get("/", response_class=HTMLResponse)
def home(request: Request):
    return templates.TemplateResponse("index.html", {"request": request})

@app.get("/pdf_editor")
def pdf_editor(request: Request):
    return templates.TemplateResponse("pdf_editor.html", {"request": request})

@app.get("/pdf_merge")
def pdf_merge(request: Request):
    return templates.TemplateResponse("pdf_merge.html", {"request": request})

@app.get("/pdf_compress")
def pdf_compress(request: Request):
    return templates.TemplateResponse("pdf_compress.html", {"request": request})

@app.get("/pdf-to-word")
def pdf_to_word_page(request: Request):
    return templates.TemplateResponse("pdf_to_word.html", {"request": request})

@app.get("/pdf-to-ppt")
def pdf_to_ppt_page(request: Request):
    return templates.TemplateResponse("pdf_to_ppt.html", {"request": request})

@app.get("/word-to-pdf")
def word_to_pdf_page(request: Request):
    return templates.TemplateResponse("word_to_pdf.html", {"request": request})

# =========================================================
# API ROUTES
# =========================================================

# ---------- PDF EDITOR SAVE ----------
@app.post("/save-pdf")
async def save_pdf(data: dict = Body(...)):
    images = []

    for img_data in data["images"]:
        img_bytes = base64.b64decode(img_data.split(",")[1])
        img = Image.open(io.BytesIO(img_bytes)).convert("RGB")
        images.append(img)

    pdf_bytes = io.BytesIO()
    images[0].save(
        pdf_bytes,
        format="PDF",
        save_all=True,
        append_images=images[1:]
    )
    pdf_bytes.seek(0)

    return StreamingResponse(
        pdf_bytes,
        media_type="application/pdf",
        headers={"Content-Disposition": "attachment; filename=edited.pdf"}
    )


# ---------- PDF → WORD ----------
@app.post("/convert/pdf-to-word")
async def pdf_to_word(file: UploadFile = File(...)):
    input_path = os.path.join(UPLOAD_DIR, file.filename)
    output_path = os.path.join(OUTPUT_DIR, file.filename.replace(".pdf", ".docx"))

    with open(input_path, "wb") as f:
        f.write(await file.read())

    cv = Converter(input_path)
    cv.convert(output_path)
    cv.close()

    return FileResponse(output_path, filename="converted.docx")


# ---------- WORD → PDF ----------
@app.post("/convert/word-to-pdf")
async def word_to_pdf(file: UploadFile = File(...)):
    input_path = os.path.join(UPLOAD_DIR, file.filename)

    with open(input_path, "wb") as f:
        f.write(await file.read())

    subprocess.run([
        LIBREOFFICE_PATH,
        "--headless",
        "--convert-to", "pdf",
        "--outdir", OUTPUT_DIR,
        input_path
    ])

    output_file = os.path.join(
        OUTPUT_DIR,
        os.path.splitext(file.filename)[0] + ".pdf"
    )

    if not os.path.exists(output_file):
        raise HTTPException(status_code=500, detail="Conversion failed")

    return FileResponse(output_file, filename="converted.pdf")


# ---------- PDF → PPT ----------
@app.post("/convert/pdf-to-ppt")
async def pdf_to_ppt(file: UploadFile = File(...)):
    try:
        input_path = os.path.join(UPLOAD_DIR, file.filename)

        with open(input_path, "wb") as f:
            f.write(await file.read())

        images = convert_from_path(
            input_path,
            poppler_path=r"C:\Users\yveer\Downloads\Release-25.12.0-0\poppler-25.12.0\Library\bin"
        )

        prs = Presentation()

        for i, img in enumerate(images):
            img_path = os.path.join(OUTPUT_DIR, f"page_{i}.jpg")
            img.save(img_path, "JPEG")

            slide = prs.slides.add_slide(prs.slide_layouts[6])

            slide.shapes.add_picture(
                img_path,
                Inches(0),
                Inches(0),
                width=prs.slide_width,
                height=prs.slide_height
            )

        output_path = os.path.join(
            OUTPUT_DIR,
            file.filename.replace(".pdf", ".pptx")
        )

        prs.save(output_path)

        return FileResponse(output_path, filename="converted.pptx")

    except Exception as e:
        print("ERROR:", str(e))
        return {"error": str(e)}